from pypdf import PdfReader
from unstructured.partition.pdf import partition_pdf
from docx import Document
import pdfplumber
import fitz  # PyMuPDF
import logging
import os
from datetime import datetime
import json
from openpyxl import load_workbook
import pandas as pd
from pptx import Presentation



logger = logging.getLogger(__name__)
"""
PDF文档加载服务类
    这个服务类提供了多种PDF文档加载方法，支持不同的加载策略和分块选项。
    主要功能：
    1. 支持多种PDF解析库：
        - PyMuPDF (fitz): 适合快速处理大量PDF文件，性能最佳
        - PyPDF: 适合简单的PDF文本提取，依赖较少
        - pdfplumber: 适合需要处理表格或需要文本位置信息的场景
        - unstructured: 适合需要更好的文档结构识别和灵活分块策略的场景
    
    2. 文档加载特性：
        - 保持页码信息
        - 支持文本分块
        - 提供元数据存储
        - 支持不同的加载策略（使用unstructured时）
 """
class LoadingService:
    """
    文档加载服务类，提供多种PDF文档加载和处理方法。
    
    属性:
        total_pages (int): 当前加载PDF文档的总页数
        current_page_map (list): 存储当前文档的页面映射信息，每个元素包含页面文本和页码
    """
    
    def __init__(self):
        self.total_pages = 0
        self.current_page_map = []
    
    def load_pdf(self, file_path: str, method: str, strategy: str = None, chunking_strategy: str = None, chunking_options: dict = None) -> str:
        """
        加载PDF文档的主方法，支持多种加载策略。

        参数:
            file_path (str): 文件路径
            method (str): 加载方法，支持 'pymupdf', 'pypdf', 'pdfplumber', 'unstructured', 'text', 'markdown', 'word', 'excel', 'csv', 'ppt'
            strategy (str, optional): 使用unstructured方法时的策略，可选 'fast', 'hi_res', 'ocr_only'
            chunking_strategy (str, optional): 文本分块策略，可选 'basic', 'by_title'
            chunking_options (dict, optional): 分块选项配置

        返回:
            str: 提取的文本内容
        """
        try:
            if method == "pymupdf":
                return self._load_with_pymupdf(file_path)
            elif method == "pypdf":
                return self._load_with_pypdf(file_path)
            elif method == "pdfplumber":
                return self._load_with_pdfplumber(file_path)
            elif method == "unstructured":
                return self._load_with_unstructured(
                    file_path, 
                    strategy=strategy,
                    chunking_strategy=chunking_strategy,
                    chunking_options=chunking_options
                )
            elif method == "text":
                return self._load_with_text(file_path)
            elif method == "markdown":
                return self._load_with_markdown(file_path)
            elif method == "word":
                return self._load_with_word(file_path)
            elif method == "excel":
                return self._load_with_excel(file_path)
            elif method == "csv":
                return self._load_with_csv(file_path)
            elif method == "ppt":
                return self._load_with_ppt(file_path)
            else:
                raise ValueError(f"Unsupported loading method: {method}")
        except Exception as e:
            logger.error(f"Error loading file with {method}: {str(e)}")
            raise
    
    def get_total_pages(self) -> int:
        """
        获取当前加载文档的总页数。

        返回:
            int: 文档总页数
        """
        return max(page_data['page'] for page_data in self.current_page_map) if self.current_page_map else 0
    
    def get_page_map(self) -> list:
        """
        获取当前文档的页面映射信息。

        返回:
            list: 包含每页文本内容和页码的列表
        """
        return self.current_page_map
    
    def _load_with_pymupdf(self, file_path: str) -> str:
        """
        使用PyMuPDF库加载PDF文档。
        适合快速处理大量PDF文件，性能最佳。

        参数:
            file_path (str): PDF文件路径

        返回:
            str: 提取的文本内容
        """
        text_blocks = []
        try:
            with fitz.open(file_path) as doc:
                self.total_pages = len(doc)
                for page_num, page in enumerate(doc, 1):
                    text = page.get_text("text")
                    if text.strip():
                        text_blocks.append({
                            "text": text.strip(),
                            "page": page_num
                        })
            self.current_page_map = text_blocks
            return "\n".join(block["text"] for block in text_blocks)
        except Exception as e:
            logger.error(f"PyMuPDF error: {str(e)}")
            raise
    
    def _load_with_pypdf(self, file_path: str) -> str:
        """
        使用PyPDF库加载PDF文档。
        适合简单的PDF文本提取，依赖较少。

        参数:
            file_path (str): PDF文件路径

        返回:
            str: 提取的文本内容
        """
        try:
            text_blocks = []
            with open(file_path, "rb") as file:
                pdf = PdfReader(file)
                self.total_pages = len(pdf.pages)
                for page_num, page in enumerate(pdf.pages, 1):
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        text_blocks.append({
                            "text": page_text.strip(),
                            "page": page_num
                        })
            self.current_page_map = text_blocks
            return "\n".join(block["text"] for block in text_blocks)
        except Exception as e:
            logger.error(f"PyPDF error: {str(e)}")
            raise
    
    def _load_with_unstructured(self, file_path: str, strategy: str = "fast", chunking_strategy: str = "basic", chunking_options: dict = None) -> str:
        """
        使用unstructured库加载PDF文档。
        适合需要更好的文档结构识别和灵活分块策略的场景。

        参数:
            file_path (str): PDF文件路径
            strategy (str): 加载策略，默认'fast'
            chunking_strategy (str): 分块策略，默认'basic'
            chunking_options (dict): 分块选项配置

        返回:
            str: 提取的文本内容
        """
        try:
            strategy_params = {
                "fast": {"strategy": "fast"},
                "hi_res": {"strategy": "hi_res"},
                "ocr_only": {"strategy": "ocr_only"}
            }            
         
            # Prepare chunking parameters based on strategy
            chunking_params = {}
            if chunking_strategy == "basic":
                chunking_params = {
                    "max_characters": chunking_options.get("maxCharacters", 4000),
                    "new_after_n_chars": chunking_options.get("newAfterNChars", 3000),
                    "combine_text_under_n_chars": chunking_options.get("combineTextUnderNChars", 2000),
                    "overlap": chunking_options.get("overlap", 200),
                    "overlap_all": chunking_options.get("overlapAll", False)
                }
            elif chunking_strategy == "by_title":
                chunking_params = {
                    "chunking_strategy": "by_title",
                    "combine_text_under_n_chars": chunking_options.get("combineTextUnderNChars", 2000),
                    "multipage_sections": chunking_options.get("multiPageSections", False)
                }
            
            # Combine strategy parameters with chunking parameters
            params = {**strategy_params.get(strategy, {"strategy": "fast"}), **chunking_params}
            
            elements = partition_pdf(file_path, **params)
            
            # Add debug logging
            for elem in elements:
                logger.debug(f"Element type: {type(elem)}")
                logger.debug(f"Element content: {str(elem)}")
                logger.debug(f"Element dir: {dir(elem)}")
            
            text_blocks = []
            pages = set()
            
            for elem in elements:
                metadata = elem.metadata.__dict__
                page_number = metadata.get('page_number')
                
                if page_number is not None:
                    pages.add(page_number)
                    
                    # Convert element to a serializable format
                    cleaned_metadata = {}
                    for key, value in metadata.items():
                        if key == '_known_field_names':
                            continue
                        
                        try:
                            # Try JSON serialization to test if value is serializable
                            json.dumps({key: value})
                            cleaned_metadata[key] = value
                        except (TypeError, OverflowError):
                            # If not serializable, convert to string
                            cleaned_metadata[key] = str(value)
                    
                    # Add additional element information
                    cleaned_metadata['element_type'] = elem.__class__.__name__
                    cleaned_metadata['id'] = str(getattr(elem, 'id', None))
                    cleaned_metadata['category'] = str(getattr(elem, 'category', None))
                    
                    text_blocks.append({
                        "text": str(elem),
                        "page": page_number,
                        "metadata": cleaned_metadata
                    })
            
            self.total_pages = max(pages) if pages else 0
            self.current_page_map = text_blocks
            return "\n".join(block["text"] for block in text_blocks)
            
        except Exception as e:
            logger.error(f"Unstructured error: {str(e)}")
            raise
    
    def _load_with_pdfplumber(self, file_path: str) -> str:
        """
        使用pdfplumber库加载PDF文档。
        适合需要处理表格或需要文本位置信息的场景。

        参数:
            file_path (str): PDF文件路径

        返回:
            str: 提取的文本内容
        """
        text_blocks = []
        try:
            with pdfplumber.open(file_path) as pdf:
                self.total_pages = len(pdf.pages)
                for page_num, page in enumerate(pdf.pages, 1):
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        text_blocks.append({
                            "text": page_text.strip(),
                            "page": page_num
                        })
            self.current_page_map = text_blocks
            return "\n".join(block["text"] for block in text_blocks)
        except Exception as e:
            logger.error(f"pdfplumber error: {str(e)}")
            raise
    
    def _load_with_text(self, file_path: str) -> str:
        """
        使用纯文本方式加载 .txt 文件。
        """
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
            text_blocks = [{"text": content, "page": 1}]
            self.total_pages = 1
            self.current_page_map = text_blocks
            return content
        except Exception as e:
            logger.error(f"Text loading error: {str(e)}")
            raise
    
    def _load_with_markdown(self, file_path: str) -> str:
        """
        使用纯文本方式加载 .md 文件（复用 _load_with_text）。
        """
        return self._load_with_text(file_path)
    
    def _load_with_word(self, file_path: str) -> str:
        """
        使用 unstructured 加载 Word (.docx) 文件。
        """
        try:
            doc = Document(file_path)
            text_blocks = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    text_blocks.append({"text": text, "page": 1})
            self.total_pages = 1
            self.current_page_map = text_blocks
            return "\n".join(block["text"] for block in text_blocks)
        except ImportError:
            logger.error("python-docx 未安装，无法加载 Word 文件")
            raise
        except Exception as e:
            logger.error(f"Word loading error: {str(e)}")
            raise
    
    def _load_with_excel(self, file_path: str) -> str:
        """
        使用 openpyxl 加载 Excel (.xlsx/.xls) 文件。
        """
        try:
            from openpyxl import load_workbook

            workbook = load_workbook(file_path, read_only=True, data_only=True)
            text_blocks = []
            for idx, sheet in enumerate(workbook, 1):
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    rows.append("\t".join(row_values))
                sheet_text = "\n".join(rows).strip()
                if sheet_text:
                    text_blocks.append({
                        "text": sheet_text,
                        "page": idx,
                        "metadata": {"sheet_name": sheet.title}
                    })
            self.total_pages = len(text_blocks)
            self.current_page_map = text_blocks
            return "\n".join(block["text"] for block in text_blocks)
        except ImportError:
            logger.error("openpyxl 未安装，无法加载 Excel 文件")
            raise
        except Exception as e:
            logger.error(f"Excel loading error: {str(e)}")
            raise
    
    def _load_with_csv(self, file_path: str) -> str:
        """
        使用 pandas 加载 CSV 文件。
        支持自动检测编码和分隔符。
        
        参数:
            file_path (str): CSV 文件路径
            
        返回:
            str: 提取的文本内容
        """
        try:
            # 尝试自动检测编码
            encodings = ['utf-8', 'gbk', 'gb2312', 'iso-8859-1']
            df = None
            detected_encoding = None
            
            for encoding in encodings:
                try:
                    df = pd.read_csv(file_path, encoding=encoding)
                    detected_encoding = encoding
                    break
                except UnicodeDecodeError:
                    continue
                except Exception as e:
                    logger.warning(f"尝试使用 {encoding} 编码读取CSV失败: {str(e)}")
                    continue
            
            if df is None:
                raise ValueError("无法使用已知编码读取CSV文件")
                
            # 将DataFrame转换为字符串格式
            text_blocks = []
            chunk_size = 1000  # 每个块的最大行数
            
            for i in range(0, len(df), chunk_size):
                chunk = df.iloc[i:i + chunk_size]
                chunk_text = chunk.to_string(index=False)
                if chunk_text.strip():
                    text_blocks.append({
                        "text": chunk_text.strip(),
                        "page": (i // chunk_size) + 1,
                        "metadata": {
                            "start_row": i,
                            "end_row": min(i + chunk_size, len(df)),
                            "total_rows": len(df),
                            "encoding": detected_encoding
                        }
                    })
            
            self.total_pages = len(text_blocks)
            self.current_page_map = text_blocks
            return "\n".join(block["text"] for block in text_blocks)
            
        except ImportError:
            logger.error("pandas 未安装，无法加载 CSV 文件")
            raise
        except Exception as e:
            logger.error(f"CSV loading error: {str(e)}")
            raise
    
    def _load_with_ppt(self, file_path: str) -> str:
        """
        使用 python-pptx 加载 PowerPoint (.pptx) 文件。
        提取幻灯片中的文本内容，包括标题、正文、表格、备注等。
        
        参数:
            file_path (str): PPT 文件路径
            
        返回:
            str: 提取的文本内容
        """
        try:
            prs = Presentation(file_path)
            text_blocks = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []
                
                # 提取幻灯片布局名称
                layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown Layout"
                
                # 提取标题
                if slide.shapes.title:
                    title_text = slide.shapes.title.text.strip()
                    if title_text:
                        slide_texts.append(f"标题: {title_text}")
                
                # 遍历所有形状
                for shape in slide.shapes:
                    # 提取文本框内容
                    if hasattr(shape, "text") and shape != slide.shapes.title:
                        text = shape.text.strip()
                        if text:
                            slide_texts.append(text)
                    
                    # 提取表格内容
                    if shape.has_table:
                        table_rows = []
                        for row in shape.table.rows:
                            row_texts = []
                            for cell in row.cells:
                                cell_text = cell.text.strip()
                                row_texts.append(cell_text if cell_text else "")
                            if any(row_texts):  # 只添加非空行
                                table_rows.append("\t".join(row_texts))
                        if table_rows:
                            slide_texts.append("\n".join(table_rows))
                
                # 提取备注
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_texts.append(f"备注: {notes}")
                
                # 如果幻灯片有内容，添加到文本块
                if slide_texts:
                    text_blocks.append({
                        "text": "\n".join(slide_texts),
                        "page": slide_num,
                        "metadata": {
                            "layout": layout_name,
                            "has_notes": bool(slide.has_notes_slide),
                            "total_shapes": len(slide.shapes)
                        }
                    })
            
            self.total_pages = len(prs.slides)
            self.current_page_map = text_blocks
            return "\n".join(block["text"] for block in text_blocks)
            
        except ImportError:
            logger.error("python-pptx 未安装，无法加载 PPT 文件")
            raise
        except Exception as e:
            logger.error(f"PPT loading error: {str(e)}")
            raise
    
    def save_document(self, filename: str, chunks: list, metadata: dict, loading_method: str, strategy: str = None, chunking_strategy: str = None) -> str:
        """
        保存处理后的文档数据。

        参数:
            filename (str): 原PDF文件名
            chunks (list): 文档分块列表
            metadata (dict): 文档元数据
            loading_method (str): 使用的加载方法
            strategy (str, optional): 使用的加载策略
            chunking_strategy (str, optional): 使用的分块策略

        返回:
            str: 保存的文件路径
        """
        try:
            timestamp = datetime.now().strftime('%Y%m%d%H%M%S')
            base_name = filename.replace('.pdf', '').split('_')[0]
            
            # Adjust the document name to include strategy if unstructured
            if loading_method == "unstructured" and strategy:
                doc_name = f"{base_name}_{loading_method}_{strategy}_{chunking_strategy}_{timestamp}"
            else:
                doc_name = f"{base_name}_{loading_method}_{timestamp}"
            
            # 构建文档数据结构，确保所有值都是可序列化的
            document_data = {
                "filename": str(filename),
                "total_chunks": int(len(chunks)),
                "total_pages": int(metadata.get("total_pages", 1)),
                "loading_method": str(loading_method),
                "loading_strategy": str(strategy) if loading_method == "unstructured" and strategy else None,
                "chunking_strategy": str(chunking_strategy) if loading_method == "unstructured" and chunking_strategy else None,
                "chunking_method": "loaded",
                "timestamp": datetime.now().isoformat(),
                "chunks": chunks
            }
            
            # 保存到文件
            filepath = os.path.join("01-loaded-docs", f"{doc_name}.json")
            os.makedirs("01-loaded-docs", exist_ok=True)
            
            with open(filepath, 'w', encoding='utf-8') as f:
                json.dump(document_data, f, ensure_ascii=False, indent=2)
                
            return filepath
            
        except Exception as e:
            logger.error(f"Error saving document: {str(e)}")
            raise

if __name__ == "__main__":
    loading_service = LoadingService()
    loading_service.load_pdf("./services/test.pptx", "ppt")
    print(loading_service.get_page_map())
